from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION (Identique HTML) ---
COLORS = {
    "bg": RGBColor(250, 250, 250),           # #FAFAFA
    "text_main": RGBColor(10, 10, 10),       # #0A0A0A
    "text_muted": RGBColor(106, 106, 106),   # #6A6A6A
    "blue": RGBColor(37, 99, 235),           # #2563EB (text-blue-600)
    "purple": RGBColor(147, 51, 234),        # #9333EA (text-purple-600)
    "card_bg": RGBColor(255, 255, 255),      # White
    "border": RGBColor(229, 229, 229)        # #E5E5E5
}

def create_slide(prs):
    # Slide vide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    # Fond
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]
    return slide

def add_centered_text(slide, text, size, color, bold=False, y_pos=Inches(1), font_name='Arial'):
    # Boite de texte largeur page pour centrage parfait
    textbox = slide.shapes.add_textbox(Inches(0), y_pos, prs.slide_width, Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return textbox

def create_card_centered(slide, y_start, height, border_color_rgb, icon_bg_rgb, icon_char, title, subtitle, tag):
    # Dimensions Card
    card_w = Inches(7) 
    card_x = (prs.slide_width - card_w) / 2
    
    # 1. Fond Carte (Glassmorphism simulé: Blanc + Bordure fine)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y_start, card_w, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card_bg"]
    card.line.color.rgb = border_color_rgb
    card.line.width = Pt(3) # Bordure épaisse colorée comme dans HTML (border-l-4 / border-t-4)

    # 2. Contenu Centré
    # Icone Rond
    icon_size = Inches(0.8)
    icon_x = card_x + (card_w - icon_size) / 2
    icon_y = y_start + Inches(0.3)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_x, icon_y, icon_size, icon_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_bg_rgb
    circle.line.fill.background()
    
    # Texte Icone (SVG simulé par emoji/caractère)
    tb = slide.shapes.add_textbox(icon_x, icon_y, icon_size, icon_size)
    p = tb.text_frame.add_paragraph()
    p.text = icon_char
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    # Hack pour centrer verticalement un peu
    tb.text_frame.margin_top = Pt(10)

    # Titre & Sous-titre
    text_y = icon_y + icon_size + Inches(0.1)
    tb_text = slide.shapes.add_textbox(card_x, text_y, card_w, Inches(1))
    
    p = tb_text.text_frame.add_paragraph()
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.color.rgb = COLORS["text_main"] # Slate-800
    p.font.size = Pt(18)
    
    p = tb_text.text_frame.add_paragraph()
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = COLORS["text_muted"] # Slate-500
    p.font.size = Pt(12)
    p.space_before = Pt(5)

    # Tag (Gris en bas)
    tag_w = Inches(2)
    tag_h = Inches(0.4)
    tag_x = card_x + (card_w - tag_w) / 2
    tag_y = y_start + height - tag_h - Inches(0.2)
    
    tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tag_x, tag_y, tag_w, tag_h)
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate-100
    tag_shape.line.fill.background() # Pas de bordure
    
    tag_tb = slide.shapes.add_textbox(tag_x, tag_y, tag_w, tag_h)
    p = tag_tb.text_frame.add_paragraph()
    p.text = tag
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_muted"]
    # Centrage vertical approx
    tag_tb.text_frame.margin_top = Pt(5)


def create_list_card_centered(slide, y_start, height, title, items):
    # Dimensions
    card_w = Inches(8)
    card_x = (prs.slide_width - card_w) / 2
    
    # Fond
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y_start, card_w, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card_bg"] # White/Glass
    card.line.color.rgb = COLORS["border"]
    card.line.width = Pt(1)
    
    # Titre Section
    tb = slide.shapes.add_textbox(card_x, y_start + Inches(0.2), card_w, Inches(0.5))
    p = tb.text_frame.add_paragraph()
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["text_main"] # Slate-800

    # Liste
    item_y = y_start + Inches(0.8)
    for txt in items:
        # Ligne item
        tb_item = slide.shapes.add_textbox(card_x + Inches(1), item_y, card_w - Inches(2), Inches(0.4))
        p = tb_item.text_frame.add_paragraph()
        p.text = txt
        p.alignment = PP_ALIGN.CENTER # Tout centrer comme demandé
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_muted"] # Slate-600
        
        # Fond item subtil (bg-white/50) -> impossible en pptx siple, on laisse transparent
        # On ajoute juste un petit bullet visuel ou cadre si on voulait pousser, mais text simple suffit
        
        item_y += Inches(0.5)

# --- MAIN ---
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9
prs.slide_height = Inches(7.5)

# SLIDE 1 : TITRE
s1 = create_slide(prs)

# Icone Microscope (Simulé)
box = s1.shapes.add_textbox(Inches(0), Inches(1.5), prs.slide_width, Inches(1.5))
p = box.text_frame.add_paragraph()
p.text = "🔬"
p.font.size = Pt(60)
p.alignment = PP_ALIGN.CENTER

# Titre
add_centered_text(s1, "RÉSEAU SCIENTIFIQUE\nScientifique", 60, COLORS["text_main"], bold=True, y_pos=Inches(3))

# Gradient Line (Simulée par rectangle bleu)
line = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, (prs.slide_width - Inches(3))/2, Inches(5.2), Inches(3), Inches(0.1))
line.fill.solid()
line.fill.fore_color.rgb = COLORS["blue"] # On met bleu, impossible de faire gradient css facilement
line.line.fill.background()

# Sous-titre
add_centered_text(s1, "PROJET GRAPHE & OPEN DATA - L3 MIASHS", 14, COLORS["text_muted"], y_pos=Inches(5.5))

# Auteur (Badge)
badge_w = Inches(3)
badge_x = (prs.slide_width - badge_w) / 2
badge_shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, Inches(6.5), badge_w, Inches(0.8))
badge_shape.fill.solid()
badge_shape.fill.fore_color.rgb = COLORS["card_bg"]
badge_shape.line.color.rgb = COLORS["border"]

tb = s1.shapes.add_textbox(badge_x, Inches(6.6), badge_w, Inches(0.5))
p = tb.text_frame.add_paragraph()
p.text = "Réalisé par\nTerrel NUENTSA"
p.font.size = Pt(12)
p.font.color.rgb = COLORS["text_main"]
p.alignment = PP_ALIGN.CENTER


# SLIDE 2 : CONCEPT
s2 = create_slide(prs)

# Titre
add_centered_text(s2, "Le Concept", 44, COLORS["text_main"], bold=True, y_pos=Inches(0.5))
# Intro Text
add_centered_text(s2, "L'histoire des sciences n'est pas une ligne droite,\nmais un tissu complexe d'influences interconnectées.", 18, COLORS["text_muted"], y_pos=Inches(1.5))

# Card 1: Entités (Bleu)
create_card_centered(s2, Inches(3), Inches(3.5), COLORS["blue"], RGBColor(219, 234, 254), "🧬", 
                     "Les Entités", "Scientifiques, mathématiciens & philosophes.", "Nœuds (Nodes)")

# Comme on doit EMPILER (stack) pour respecter le "exactement comme html centré verticalement",
# Mais en PPTX on a moins de hauteur de scroll. On va mettre Côte à côte car stacker 2 grosses cartes dépasse la slide ?
# L'HTML utilisait "flex-col items-center" -> donc empilé.
# Mais slide PPTX = 7.5 pouces de haut.
# Titre (1.5) + Text (1) + Card1 (3.5) + Card2 (3.5) = 9.5 pouces -> CA DEPASSE.
# L'HTML scrollait ("overflow-y: scroll"). PPTX ne scrolle pas.
# SOLUTION : Je vais les mettre côte à côte mais BIEN CENTRÉES pour garder l'équilibre.
# L'utilisateur a dit "exactement la meme", mais si ça rentre pas... je vais adapter pour "l'esprit" centré.
# ATTENTION : Si l'utilisateur veut "exactement", il veut peut être que je mette que les éléments principaux.
# Je vais mettre côte à côte, c'est le seul moyen PPTX propre.

# Modif fonction create_card pour x
def create_card_positioned(slide, x, y, w, h, border_color, icon_bg, icon, title, subtitle, tag):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card_bg"]
    card.line.color.rgb = border_color
    card.line.width = Pt(3)
    
    # Icone
    icon_size = Inches(0.6)
    icon_x = x + (w - icon_size)/2
    icon_y = y + Inches(0.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_x, icon_y, icon_size, icon_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = icon_bg
    circle.line.fill.background()
    
    tb = slide.shapes.add_textbox(icon_x, icon_y, icon_size, icon_size)
    p = tb.text_frame.add_paragraph()
    p.text = icon
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)

    # Texte
    tb = slide.shapes.add_textbox(x, y + Inches(1), w, Inches(1.5))
    p = tb.text_frame.add_paragraph()
    p.text = title
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    
    p = tb.text_frame.add_paragraph()
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["text_muted"]

    # Tag
    tag_w = Inches(1.5)
    ts = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + (w-tag_w)/2, y + h - Inches(0.6), tag_w, Inches(0.4))
    ts.fill.solid()
    ts.fill.fore_color.rgb = RGBColor(241, 245, 249)
    ts.line.fill.background()
    
    tb2 = slide.shapes.add_textbox(x + (w-tag_w)/2, y + h - Inches(0.55), tag_w, Inches(0.4))
    p = tb2.text_frame.add_paragraph()
    p.text = tag
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_muted"]

# Card 1 (Gauche, mais centrée bloc)
# Largeur slide 13.33. Card width 4. Gap 1. Total 9.
create_card_positioned(s2, Inches(2.5), Inches(3.5), Inches(4), Inches(3), COLORS["blue"], RGBColor(219, 234, 254), "🧬", "Les Entités", "Scientifiques,\nPhilosophes.", "Nœuds")
# Card 2 (Droite)
create_card_positioned(s2, Inches(7), Inches(3.5), Inches(4), Inches(3), COLORS["purple"], RGBColor(233, 213, 255), "🔗", "Les Relations", "Influences\nBio-sourcées.", "Arêtes")


# SLIDE 3 : DONNÉES
s3 = create_slide(prs)
add_centered_text(s3, "Analyse & Données", 44, COLORS["text_main"], bold=True, y_pos=Inches(0.5))

# Données (Box 1)
# On réutilise create_card_positioned mais pour liste
def create_list_box(slide, x, y, w, h, title, items):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["card_bg"]
    shape.line.color.rgb = COLORS["border"]
    
    tb = slide.shapes.add_textbox(x, y + Inches(0.2), w, Inches(0.5))
    p = tb.text_frame.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    iy = y + Inches(0.8)
    for it in items:
        tb_i = slide.shapes.add_textbox(x, iy, w, Inches(0.4))
        p = tb_i.text_frame.add_paragraph()
        p.text = it
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["text_muted"]
        iy += Inches(0.4)

create_list_box(s3, Inches(2.5), Inches(2), Inches(4), Inches(4.5), "📊 Données", 
                ["Wikipédia API", "LLM Groq/Llama3", "~250 Nœuds", "JSON Structuré"])

create_list_box(s3, Inches(7), Inches(2), Inches(4), Inches(4.5), "🧐 Questions", 
                ["1. Qui influence qui ?", "2. Cercles Scientifiques ?", "3. Chemins d'Idées ?"])


output_path = "output/presentation_final.pptx"
prs.save(output_path)
print(f"✅ Slide Finale générée : {output_path}")

def main():
    # --- MAIN ---
    global prs # Need to define prs if used in functions that don't pass it, but helper functions use 'slide' or 'prs' passed?
    # Helper functions add_centered_text uses 'prs.slide_width' which assumes 'prs' is global.
    # To be safe and quick, I will just remove the if __name__ block and let the script run linearly, OR define create_presentation properly.
    
    # Actually, looking at lines 28-38: add_centered_text uses 'prs.slide_width'. 'prs' is defined at line 149.
    # If I wrap line 149+ in a function, 'add_centered_text' will fail unless 'prs' is passed or global.
    # The clean fix given I can't refactor everything easily:
    # Just remove the 'if __name__' block and the call. The script executes linearly from line 149.
    pass

if __name__ == "__main__":
    pass
