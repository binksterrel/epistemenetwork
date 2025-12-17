from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

# Palette de couleurs "Premium Science"
COLORS = {
    "bg": "FAFAFA",           # Blanc cassé fond
    "primary": "2563EB",      # Bleu Roi
    "secondary": "7C3AED",    # Violet Profond
    "accent": "06B6D4",       # Cyan
    "text_dark": "111827",    # Presque noir
    "text_gray": "4B5563",    # Gris texte
    "card_bg": "FFFFFF"       # Blanc pur
}

def apply_background_style(slide):
    # Ajoute une bande décorative latérale
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(COLORS["primary"])
    shape.line.fill.background()
    
    # Ajoute un petit accent en bas à droite
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9), Inches(6.5), Inches(1), Inches(1)
    )
    shape.rotation = 180
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(COLORS["secondary"])
    shape.line.fill.background()

def create_card(slide, x, y, w, h, title_text, body_text, icon_char, border_color_hex):
    # Ombre (simulée par un rectangle gris décalé)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.05), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shadow.line.fill.background()
    
    # Carte principale
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(COLORS["card_bg"])
    card.line.color.rgb = hex_to_rgb(border_color_hex)
    card.line.width = Pt(2)

    # Icône
    textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(0.8))
    p = textbox.text_frame.add_paragraph()
    p.text = icon_char
    p.font.size = Pt(40)
    
    # Titre Carte
    textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.5))
    p = textbox.text_frame.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(COLORS["text_dark"])
    p.font.name = 'Arial'

    # Corps Carte
    textbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.4), w - Inches(0.4), h - Inches(1.6))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = body_text
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(COLORS["text_gray"])
    p.font.name = 'Arial'

def create_presentation():
    prs = Presentation()
    
    # --- SLIDE 1 : Titre Stylé ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(COLORS["bg"])
    bg.line.fill.background()

    # Formes Géométriques Arrière-plan
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(1), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
    circle.fill.transparency = 0.8
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = hex_to_rgb(COLORS["secondary"])
    circle2.fill.transparency = 0.9
    circle2.line.fill.background()

    # Titre Principal
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p = textbox.text_frame.add_paragraph()
    p.text = "RÉSEAU SCIENTIFIQUE\nScientifique"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.name = 'Arial Black' 
    p.font.color.rgb = hex_to_rgb(COLORS["text_dark"])
    p.alignment = PP_ALIGN.CENTER
    
    # Ligne de séparation
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.8), Inches(2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(COLORS["primary"])
    line.line.fill.background()

    # Sous-titre
    textbox = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = "PROJET GRAPHE & OPEN DATA"
    p.font.size = Pt(16)
    p.font.name = 'Arial'
    p.font.color.rgb = hex_to_rgb(COLORS["text_gray"])
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

    # Auteur Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.2), Inches(3), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = hex_to_rgb(COLORS["card_bg"])
    badge.line.color.rgb = hex_to_rgb("E5E7EB")
    
    textbox = slide.shapes.add_textbox(Inches(3.5), Inches(6.35), Inches(3), Inches(0.5))
    p = textbox.text_frame.add_paragraph()
    p.text = "Réalisé par Terrel NUENTSA"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.name = 'Arial'
    p.font.color.rgb = hex_to_rgb(COLORS["secondary"])


    # --- SLIDE 2 : Concept (Cards Style) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background_style(slide)
    
    # Titre
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = "Le Concept"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.name = 'Arial'
    p.font.color.rgb = hex_to_rgb(COLORS["text_dark"])

    # Intro Text
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = "L'histoire des sciences n'est pas une ligne droite, mais un réseau complexe d'interconnexions."
    p.font.size = Pt(18)
    p.font.color.rgb = hex_to_rgb(COLORS["text_gray"])

    # Cards
    create_card(slide, Inches(1), Inches(3), Inches(3.8), Inches(3), 
                "Les Entités", 
                "Nous extrayons les scientifiques, mathématiciens et philosophes présents sur Wikipédia comme Nœuds du graphe.", 
                "🧠", COLORS["primary"])

    create_card(slide, Inches(5.2), Inches(3), Inches(3.8), Inches(3), 
                "Les Relations", 
                "L'IA analyse les biographies pour détecter qui a influencé qui (Relation dirigée A → B) avec précision.", 
                "🔗", COLORS["secondary"])


    # --- SLIDE 3 : Données (Infographic Style) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background_style(slide)

    # Titre
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    p = textbox.text_frame.add_paragraph()
    p.text = "Données & Enjeux"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.name = 'Arial'
    p.font.color.rgb = hex_to_rgb(COLORS["text_dark"])

    # Zone Data (Gauche)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2), Inches(4), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb("EFF6FF") # Bleu très pâle
    shape.line.fill.background()

    textbox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(3.5), Inches(4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "SOURCES DE DONNÉES"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(COLORS["primary"])
    
    data_points = [
        ("📚", "Wikipédia API"),
        ("🤖", "LLM Groq/Llama3"),
        ("🕸️", "~250 Nœuds"),
        ("💾", "JSON Structuré")
    ]
    
    for icon, text in data_points:
        p = tf.add_paragraph()
        p.text = f"\n{icon}  {text}"
        p.font.size = Pt(16)
        p.font.name = 'Arial'

    # Zone Questions (Droite)
    textbox = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.5), Inches(4))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "QUESTIONS POSÉES"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(COLORS["secondary"])
    
    questions = [
        "1. Hubs & Connecteurs",
        "Quels scientifiques relient des domaines éloignés ?",
        "2. Détection de Communautés",
        "Peut-on voir émerger des écoles de pensée ?",
        "3. Impact visuel",
        "Comment représenter le poids historique d'une idée ?"
    ]
    
    for i, q in enumerate(questions):
        p = tf.add_paragraph()
        p.text = q
        if i % 2 == 0: # Titre question
            p.font.bold = True
            p.font.size = Pt(16)
            p.space_before = Pt(12)
            p.font.color.rgb = hex_to_rgb(COLORS["text_dark"])
        else: # Détail
            p.font.size = Pt(14)
            p.font.color.rgb = hex_to_rgb(COLORS["text_gray"])


    output_path = "output/presentation_premium.pptx" # Nouveau nom
    prs.save(output_path)
    print(f"✅ Présentation Premium générée : {output_path}")

if __name__ == "__main__":
    create_presentation()
